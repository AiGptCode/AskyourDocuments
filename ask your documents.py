import os
import glob
import fitz
from PIL import Image
import textract
import pandas as pd
from pptx import Presentation
from transformers import AutoTokenizer, AutoModelForQuestionAnswering
from langchain.text_splitter import CharacterTextSplitter
from langchain.embeddings.huggingface import HuggingFaceEmbeddings
from langchain.vectorstores import FAISS
from langchain.chains.question_answering import load_qa_chain
import concurrent.futures

def extract_text_from_pdf(file):
    doc = fitz.open(file)
    text = ""
    for page in doc:
        text += page.get_text()
    return text

def extract_text_from_image(file):
    text = textract.process(file, method='tesseract').decode("utf-8")
    return text

def extract_text_from_excel(file):
    df = pd.read_excel(file)
    return df.to_string(index=False, header=False)

def extract_text_from_powerpoint(file):
    prs = Presentation(file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text += shape.text_frame.text
    return text

def extract_text_from_file(file):
    return textract.process(file).decode("utf-8")

def process_document(file):
    supported_formats = [".pdf", ".txt", ".doc", ".docx", ".png", ".jpg", ".jpeg", ".xlsx", ".pptx"]
    if file.endswith(".pdf"):
        return extract_text_from_pdf(file)
    elif file.endswith((".png", ".jpg", ".jpeg")):
        return extract_text_from_image(file)
    elif file.endswith(".xlsx"):
        return extract_text_from_excel(file)
    elif file.endswith(".pptx"):
        return extract_text_from_powerpoint(file)
    else:
        return extract_text_from_file(file)

def process_directory(directory):
    documents = {}
    for filename in glob.glob(f"{directory}/**/*", recursive=True):
        try:
            text = process_document(filename)
            documents[filename] = text
        except Exception as e:
            print(f"Error processing document {filename}: {e}")
    return documents

def handle_question(knowledge_base, question):
    model_name = "Helsinki-NLP/mistral-7b"
    tokenizer = AutoTokenizer.from_pretrained(model_name)
    model = AutoModelForQuestionAnswering.from_pretrained(model_name)
    chain = load_qa_chain(model, chain_type="stuff")
    responses = []
    for filename, docs in knowledge_base.items():
        similar_docs = docs.similarity_search(question)
        response = chain.run(input_documents=similar_docs, question=question)
        responses.append(f"{filename}: {response}")
    return responses

def main():
    api_key = input("Enter your Hugging Face API key: ")
    os.environ["HUGGINGFACE_API_TOKEN"] = api_key
    directory = input("Enter the directory path containing your documents: ")
    if os.path.isdir(directory):
        documents = process_directory(directory)
        text_splitter = CharacterTextSplitter(
            separator="\n",
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len
        )
        chunks = {}
        for filename, text in documents.items():
            chunks[filename] = text_splitter.split_text(text)
        embeddings = HuggingFaceEmbeddings()
        with concurrent.futures.ProcessPoolExecutor() as executor:
            futures = {executor.submit(FAISS.from_texts, document_chunks, embeddings): filename for filename, document_chunks in chunks.items()}
            knowledge_base = {filename: future.result() for filename, future in futures.items()}
        while True:
            user_question = input("Ask a question about your documents (or type 'quit' to exit): ")
            if user_question.lower() == "quit":
                break
            responses = handle_question(knowledge_base, user_question)
            for response in responses:
                print(response)
    else:
        print("The provided path is not a valid directory.")

if __name__ == "__main__":
    main()
